#!/usr/bin/env python3
"""Create a PowerPoint presentation about Claude Skills MCP architecture."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create the presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Color scheme
    DARK_BLUE = RGBColor(28, 40, 51)      # #1C2833
    TEAL = RGBColor(94, 168, 167)         # #5EA8A7
    CORAL = RGBColor(254, 68, 71)         # #FE4447
    GRAY = RGBColor(46, 64, 83)           # #2E4053
    LIGHT_GRAY = RGBColor(244, 246, 246)  # #F4F6F6
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(4.4), Inches(1.8),
        Inches(1.2), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = CORAL
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Claude Skills MCP Server"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.2),
        Inches(8), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI Agent Skills를 위한 지능형 검색 시스템"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title with underline
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "프로젝트 개요"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0.03))
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL
    underline.line.fill.background()
    
    # Content bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    bullets = [
        "Anthropic의 Agent Skills 프레임워크를 MCP로 확장",
        "벡터 임베딩 기반 시맨틱 검색",
        "123개 이상의 큐레이션된 스킬 제공",
        "GitHub 저장소와 로컬 디렉토리 지원",
        "API 키 불필요, 완전 로컬 동작"
    ]
    
    for bullet_text in bullets:
        p = content_frame.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_before = Pt(10)
    
    # Highlight box
    highlight = slide.shapes.add_shape(1, Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = TEAL
    highlight.line.fill.background()
    
    highlight_text = highlight.text_frame
    highlight_text.text = "핵심 가치: Anthropic의 Skills 시스템을 Cursor, Codex, GPT-5 등\n모든 AI 모델에서 사용 가능하게 만듭니다."
    for p in highlight_text.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Two-Package Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Two-Package 아키텍처"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0.03))
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL
    underline.line.fill.background()
    
    # Frontend component
    frontend = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(3.5), Inches(3))
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = LIGHT_GRAY
    frontend.line.color.rgb = TEAL
    frontend.line.width = Pt(2)
    
    frontend_text = frontend.text_frame
    frontend_text.text = "Frontend\n~15 MB\n\n• 즉시 시작 (<5초)\n• MCP 서버 (stdio)\n• 툴 스키마 즉시 반환\n• 백엔드 자동 다운로드\n\n✓ Cursor 타임아웃 해결"
    for i, p in enumerate(frontend_text.paragraphs):
        if i == 0:
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
            p.alignment = PP_ALIGN.CENTER
        elif i == 1:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = TEAL
            p.alignment = PP_ALIGN.CENTER
        else:
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.8), Inches(1), Inches(0.5))
    arrow_text = arrow_box.text_frame
    arrow_text.text = "→"
    p = arrow_text.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CORAL
    p.alignment = PP_ALIGN.CENTER
    
    # Backend component
    backend = slide.shapes.add_shape(1, Inches(5.7), Inches(1.5), Inches(3.5), Inches(3))
    backend.fill.solid()
    backend.fill.fore_color.rgb = LIGHT_GRAY
    backend.line.color.rgb = TEAL
    backend.line.width = Pt(2)
    
    backend_text = backend.text_frame
    backend_text.text = "Backend\n~250 MB\n\n• 벡터 검색 엔진\n• PyTorch + Transformers\n• MCP 서버 (HTTP)\n• 백그라운드 자동 설치\n\n✓ 강력한 시맨틱 검색"
    for i, p in enumerate(backend_text.paragraphs):
        if i == 0:
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
            p.alignment = PP_ALIGN.CENTER
        elif i == 1:
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = TEAL
            p.alignment = PP_ALIGN.CENTER
        else:
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
    
    # Slide 4: Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "워크플로우"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0.03))
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL
    underline.line.fill.background()
    
    # Workflow steps
    steps = [
        ("1", "Frontend 즉시 시작", "Cursor가 uvx로 실행하면 5초 이내에 Frontend 시작, 툴 스키마 즉시 반환"),
        ("2", "Backend 백그라운드 다운로드", "Frontend가 백그라운드에서 Backend 패키지(~250MB) 자동 다운로드 및 설치"),
        ("3", "스킬 로딩 및 인덱싱", "GitHub에서 스킬 다운로드, 벡터 임베딩 생성, 검색 인덱스 구축"),
        ("4", "시맨틱 검색 실행", "사용자 쿼리를 임베딩으로 변환, 유사도 기반 관련 스킬 검색 및 반환")
    ]
    
    y_pos = 1.3
    for num, title, desc in steps:
        # Number circle
        circle = slide.shapes.add_shape(1, Inches(0.6), Inches(y_pos), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        
        num_text = circle.text_frame
        num_text.text = num
        p = num_text.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Content box
        content = slide.shapes.add_shape(1, Inches(1.2), Inches(y_pos), Inches(8), Inches(0.65))
        content.fill.solid()
        content.fill.fore_color.rgb = LIGHT_GRAY
        content.line.fill.background()
        
        content_text = content.text_frame
        content_text.text = f"{title}\n{desc}"
        for i, p in enumerate(content_text.paragraphs):
            if i == 0:
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
            else:
                p.font.size = Pt(12)
                p.font.color.rgb = GRAY
        
        y_pos += 0.8
    
    # Highlight box
    highlight = slide.shapes.add_shape(1, Inches(1.5), Inches(4.7), Inches(7), Inches(0.5))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = CORAL
    highlight.line.fill.background()
    
    highlight_text = highlight.text_frame
    highlight_text.text = "전체 프로세스: 첫 실행 60-120초 | 이후 즉시 사용 가능"
    p = highlight_text.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    highlight_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Slide 5: MCP Tools
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "제공되는 MCP 툴"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0.03))
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL
    underline.line.fill.background()
    
    # Tool cards
    tools = [
        ("🔍", "find_helpful_skills", "작업 설명을 기반으로 관련 스킬을 시맨틱 검색으로 찾아 랭킹된 후보 반환"),
        ("📄", "read_skill_document", "스킬의 특정 파일(스크립트, 참조, 자산)을 검색. 패턴 매칭 지원"),
        ("📋", "list_skills", "로드된 모든 스킬의 전체 목록 반환 (이름, 설명, 소스, 문서 수)")
    ]
    
    x_positions = [0.8, 3.6, 6.4]
    for i, (icon, name, desc) in enumerate(tools):
        tool_box = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(1.5), Inches(2.6), Inches(3))
        tool_box.fill.solid()
        tool_box.fill.fore_color.rgb = LIGHT_GRAY
        tool_box.line.color.rgb = TEAL
        tool_box.line.width = Pt(2)
        
        tool_text = tool_box.text_frame
        tool_text.text = f"{icon}\n\n{name}\n\n{desc}"
        for j, p in enumerate(tool_text.paragraphs):
            if j == 0:  # Icon
                p.font.size = Pt(36)
                p.alignment = PP_ALIGN.CENTER
            elif j == 2:  # Name
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
                p.alignment = PP_ALIGN.CENTER
            elif j == 4:  # Description
                p.font.size = Pt(11)
                p.font.color.rgb = GRAY
    
    # Slide 6: Skills Status
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "로드된 스킬 현황"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    underline = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(0.03))
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL
    underline.line.fill.background()
    
    # Stats boxes
    stat1 = slide.shapes.add_shape(1, Inches(1.5), Inches(1.3), Inches(3), Inches(1))
    stat1.fill.solid()
    stat1.fill.fore_color.rgb = TEAL
    stat1.line.fill.background()
    
    stat1_text = stat1.text_frame
    stat1_text.text = "123\n총 로드된 스킬"
    for i, p in enumerate(stat1_text.paragraphs):
        if i == 0:
            p.font.size = Pt(48)
            p.font.bold = True
        else:
            p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    stat2 = slide.shapes.add_shape(1, Inches(5.5), Inches(1.3), Inches(3), Inches(1))
    stat2.fill.solid()
    stat2.fill.fore_color.rgb = TEAL
    stat2.line.fill.background()
    
    stat2_text = stat2.text_frame
    stat2_text.text = "2\n스킬 소스"
    for i, p in enumerate(stat2_text.paragraphs):
        if i == 0:
            p.font.size = Pt(48)
            p.font.bold = True
        else:
            p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Source cards
    sources = [
        ("Anthropic 공식 스킬", "15개", "문서 처리, 프레젠테이션, 웹 아티팩트, MCP 빌더 등"),
        ("K-Dense AI 과학 스킬", "108개", "생물정보학, 화학정보학, 딥러닝, 데이터 분석 등")
    ]
    
    x_positions = [1.2, 5.5]
    for i, (title, count, desc) in enumerate(sources):
        source_box = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(2.6), Inches(3.8), Inches(2.3))
        source_box.fill.solid()
        source_box.fill.fore_color.rgb = LIGHT_GRAY
        source_box.line.fill.background()
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(2.6), Inches(0.04), Inches(2.3))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = CORAL
        accent_bar.line.fill.background()
        
        source_text = source_box.text_frame
        source_text.text = f"\n{title}\n\n{count}\n\n{desc}"
        for j, p in enumerate(source_text.paragraphs):
            if j == 1:  # Title
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
            elif j == 3:  # Count
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = TEAL
            elif j == 5:  # Description
                p.font.size = Pt(13)
                p.font.color.rgb = GRAY
    
    # Save presentation
    output_path = "/Users/uengine/claude-skills-mcp/claude_skills_mcp_architecture.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation created successfully!")
    print(f"📄 Output: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()


